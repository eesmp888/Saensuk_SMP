import os
import json
import hashlib
import hmac
import base64
import requests
from flask import Flask, request, abort

app = Flask(__name__)

LINE_CHANNEL_SECRET = os.environ.get("LINE_CHANNEL_SECRET", "")
LINE_CHANNEL_ACCESS_TOKEN = os.environ.get("LINE_CHANNEL_ACCESS_TOKEN", "")
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")

# ─── อ่านข้อมูลจากไฟล์ใน /data ────────────────────────────────────────────
def load_data_files():
    text = ""
    data_dir = os.path.join(os.path.dirname(__file__), "data")
    if not os.path.exists(data_dir):
        return ""
    for filename in os.listdir(data_dir):
        filepath = os.path.join(data_dir, filename)
        try:
            if filename.endswith(".pptx"):
                from pptx import Presentation
                prs = Presentation(filepath)
                text += f"\n=== ข้อมูลจาก {filename} ===\n"
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"
                    if slide_text:
                        text += f"[Slide {i}]\n{slide_text}\n"
            elif filename.endswith((".xlsx", ".xls")):
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True)
                text += f"\n=== ข้อมูลจาก {filename} ===\n"
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text += f"[Sheet: {sheet}]\n"
                    for row in ws.iter_rows(values_only=True):
                        row_text = "\t".join([str(c) if c is not None else "" for c in row])
                        if row_text.strip():
                            text += row_text + "\n"
            elif filename.endswith(".docx"):
                from docx import Document
                doc = Document(filepath)
                text += f"\n=== ข้อมูลจาก {filename} ===\n"
                for para in doc.paragraphs:
                    if para.text.strip():
                        text += para.text + "\n"
            elif filename.endswith(".txt"):
                with open(filepath, "r", encoding="utf-8") as f:
                    content = f.read().strip()
                    if content and content != "data folder":
                        text += f"\n=== ข้อมูลจาก {filename} ===\n{content}\n"
        except Exception as e:
            print(f"Error reading {filename}: {e}")
    return text

# โหลดข้อมูลตอน startup
FILE_DATA = load_data_files()

STATIC_KNOWLEDGE = """
=== ข้อมูลโครงการ ===
ชื่อโครงการ: โครงการจ้างเหมาเอกชนบริหารจัดการระบบบำบัดน้ำเสียรวม เทศบาลเมืองแสนสุข จังหวัดชลบุรี
สัญญาเลขที่: 1/2569 ลงวันที่ 30 กันยายน 2568
วันเริ่มสัญญา: 1 ตุลาคม 2568 / วันสิ้นสุดสัญญา: 30 กันยายน 2569
ผู้รับจ้าง: บริษัท เจม เอ็นไวรันเมนทัล แมเนจเม้นท์ จำกัด (GEM)
มูลค่างาน: 2,615,500 บาท (รวมภาษีมูลค่าเพิ่ม)
คณะกรรมการตรวจรับพัสดุ: 1.นายชัชวาล กอหญ้ากลาง (ประธาน) 2.นางสาวเพ็ญศรี ชายชาญณรงค์ 3.นายณรงค์ศักดิ์ กลิ่นขจร

=== รายงานมีนาคม 2569 (งวดที่ 6) ===
น้ำเสียรวม: 467,240 ลบ.ม. (เหนือ 258,120 / ใต้ 209,138)
ไฟฟ้า: 69,000 kWh / น้ำประปา: 466.85 ลบ.ม. / น้ำรดต้นไม้: 4,573 ลบ.ม.
คุณภาพน้ำเหนือ: pH 6.66/6.55, Temp 27.25/27.55C, DO 2.05/3.81 (ผ่าน)
คุณภาพน้ำใต้: pH 7.15/7.14, Temp 29.16/28.35C, DO 2.91/4.15 (ผ่าน)
BOD เหนือ 45% / ใต้ 49% (ใต้ไม่ผ่านมาตรฐาน เนื่องจาก Sludge Bulking และ Clarifier ชำรุด)
ค่าจ้างงวดที่ 6: 203,621.50 + VAT = 217,875.00 บาท

=== สถานะเครื่องจักร ณ 31 มี.ค. 2569 ===
แสนสุขเหนือ (117 รายการ):
- สถานีแหลมแท่น: Pump No.1 ปกติ / No.2 ส่งซ่อม / No.3 ชำรุด
- บางแสนสาย2: Pump No.2 ปกติ / No.1,3 ชำรุด
- สุขุมวิท: Pump No.1,2,3 เครื่อง+ควบคุมชำรุด
- โรงเหนือ: Pump No.3 ปกติ / No.1,4 ชำรุดถาวร / No.2 ปั๊มปกติตู้ชำรุด
- Aerator ปกติ: No.2,3,7,9,12,14 / ชำรุด: No.1,4,5,6,8,10,11,13,15,16
- Clarifier No.1,2,8,9: ชำรุด
แสนสุขใต้ (92 รายการ):
- หาดวอนนภา: WW Pump No.1,2 ปกติ / EFF Pump No.2 ปกติ / No.1,3 เครื่องชำรุด
- โรงใต้: Aerator No.3,7 ปกติ / No.2,4,5,6,8 ชำรุดถาวร / Clarifier No.1,2 ชำรุด
"""

def get_knowledge():
    return STATIC_KNOWLEDGE + FILE_DATA

def verify_signature(body, signature):
    hash_val = hmac.new(LINE_CHANNEL_SECRET.encode("utf-8"), body, hashlib.sha256).digest()
    return hmac.compare_digest(base64.b64encode(hash_val).decode("utf-8"), signature)

def reply_message(reply_token, text):
    requests.post(
        "https://api.line.me/v2/bot/message/reply",
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {LINE_CHANNEL_ACCESS_TOKEN}"},
        json={"replyToken": reply_token, "messages": [{"type": "text", "text": text}]}
    )

def ask_groq(user_message):
    try:
        kb = get_knowledge()
        # จำกัด knowledge ไม่เกิน 6000 ตัวอักษร
        if len(kb) > 6000:
            kb = kb[:6000] + "\n...(ข้อมูลถูกตัดบางส่วน)"
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.1-8b-instant",
                "messages": [
                    {"role": "system", "content": f'ตอบเป็นภาษาไทย กระชับ ใช้ข้อมูล:\n{kb}'},
                    {"role": "user", "content": user_message}
                ],
                "temperature": 0.3,
                "max_tokens": 300
            },
            timeout=30
        )
        return response.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"Error: {str(e)}"

TRIGGER_WORDS = ["น้องแสนสุข", "แสนสุข"]

def should_reply(event):
    source_type = event.get("source", {}).get("type", "user")
    text = event["message"]["text"]
    if source_type == "user":
        return True
    if source_type in ["group", "room"]:
        mention = event["message"].get("mention", {})
        if mention.get("mentionees"):
            return True
        for word in TRIGGER_WORDS:
            if word in text:
                return True
        return False
    return True

def clean_text(event):
    text = event["message"]["text"]
    for word in TRIGGER_WORDS:
        text = text.replace(word, "").strip()
    return text if text else event["message"]["text"]

@app.route("/webhook", methods=["POST"])
def webhook():
    signature = request.headers.get("X-Line-Signature", "")
    body = request.get_data()
    if not verify_signature(body, signature):
        abort(400)
    for event in json.loads(body).get("events", []):
        if event.get("type") == "message" and event["message"].get("type") == "text":
            if should_reply(event):
                reply_message(event["replyToken"], ask_groq(clean_text(event)))
    return "OK", 200

@app.route("/", methods=["GET"])
def health():
    return "น้องแสนสุข Bot is running! 🌊", 200

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
